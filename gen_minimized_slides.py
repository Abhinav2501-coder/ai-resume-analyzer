
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colors ────────────────────────────────────────────────────────────────────
BG_COLOR      = RGBColor(0xE8, 0xE8, 0xF0)
TITLE_COLOR   = RGBColor(0x1A, 0x1A, 0x6E)
BULLET_COLOR  = RGBColor(0x1A, 0x1A, 0x6E)
ACCENT_PURPLE = RGBColor(0x6A, 0x5A, 0xCD)
LINE_COLOR    = RGBColor(0x6A, 0x5A, 0xCD)

# ── Slide data (minimized) ────────────────────────────────────────────────────
slides_data = [
    {
        "title": "KEY LEARNINGS FROM STARTUP FAILURES",
        "bullets": [
            "Plan, research & validate the market before launching.",
            "Solve real customer problems; adapt based on feedback.",
            "Financial discipline, strong leadership & skilled teams drive success.",
        ],
    },
    {
        "title": "GROWTH STRATEGIES AND ORGANIZATIONAL SUCCESS",
        "bullets": [
            "Growth expands capabilities, market reach & overall performance.",
            "Internal: innovation & employee development; External: mergers & partnerships.",
            "Strategic leadership & teamwork ensure long-term organizational success.",
        ],
    },
    {
        "title": "REVENUE MODELS AND FINANCIAL CHALLENGES",
        "bullets": [
            "Revenue models (subscription, freemium, ads) explain how startups earn.",
            "Choose the right model based on business type, market & growth strategy.",
            "Sound financial planning overcomes capital limits & high operating costs.",
        ],
    },
    {
        "title": "LEADERSHIP AND TEAM BUILDING",
        "bullets": [
            "Leadership influences individuals & teams toward shared organizational goals.",
            "Transformational leaders inspire creativity & promote long-term commitment.",
            "Trust, open communication & motivation significantly boost team performance.",
        ],
    },
    {
        "title": "STARTUP CHALLENGES AND FAILURES",
        "bullets": [
            "Startups face uncertainty, limited capital & untested business models.",
            "Failures stem from poor planning, lack of product-market fit & weak marketing.",
            "Team conflicts & rapid scaling without strong systems increase failure risk.",
        ],
    },
]

# ── Presentation ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

W = prs.slide_width
H = prs.slide_height
blank_layout = prs.slide_layouts[6]


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_decorative_accents(slide):
    """Purple rectangle blocks on the corners to mimic the triangle decorations."""
    # Top-left block
    s = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.9), Inches(2.0))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT_PURPLE; s.line.fill.background()

    # Bottom-left block
    s2 = slide.shapes.add_shape(1, Inches(0), Inches(4.1), Inches(0.7), Inches(1.525))
    s2.fill.solid(); s2.fill.fore_color.rgb = ACCENT_PURPLE; s2.line.fill.background()

    # Top-right block
    s3 = slide.shapes.add_shape(1, W - Inches(0.8), Inches(0), Inches(0.8), Inches(1.6))
    s3.fill.solid(); s3.fill.fore_color.rgb = ACCENT_PURPLE; s3.line.fill.background()

    # Bottom-right block
    s4 = slide.shapes.add_shape(1, W - Inches(1.0), Inches(4.2), Inches(1.0), Inches(1.425))
    s4.fill.solid(); s4.fill.fore_color.rgb = ACCENT_PURPLE; s4.line.fill.background()


def add_bottom_line(slide):
    shape = slide.shapes.add_shape(1, Inches(1.0), Inches(5.0), Inches(8.0), Inches(0.02))
    shape.fill.background()
    shape.line.color.rgb = LINE_COLOR
    shape.line.width = Pt(2.5)


def add_title(slide, title_text):
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.12), Inches(8.0), Inches(1.05))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.bold      = True
    run.font.size      = Pt(27)
    run.font.color.rgb = TITLE_COLOR
    run.font.name      = "Arial Black"


def add_bullets(slide, bullets):
    txBox = slide.shapes.add_textbox(Inches(1.1), Inches(1.25), Inches(7.8), Inches(3.6))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(10)
        p.space_after  = Pt(6)

        br = p.add_run()
        br.text = "\u2022  "
        br.font.size      = Pt(21)
        br.font.color.rgb = BULLET_COLOR
        br.font.name      = "Arial"

        tr = p.add_run()
        tr.text = bullet_text
        tr.font.size      = Pt(21)
        tr.font.color.rgb = BULLET_COLOR
        tr.font.name      = "Arial"


# ── Build ─────────────────────────────────────────────────────────────────────
for data in slides_data:
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_decorative_accents(slide)
    add_title(slide, data["title"])
    add_bullets(slide, data["bullets"])
    add_bottom_line(slide)

# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.expanduser("~"), "Downloads", "minimized_slides.pptx")
prs.save(out)
print(f"Saved -> {out}")
